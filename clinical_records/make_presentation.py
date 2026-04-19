from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Color palette ─────────────────────────────────────────────────────────────
DARK_BG     = RGBColor(0x0D, 0x11, 0x17)
CARD_BG     = RGBColor(0x16, 0x1B, 0x22)
BLUE        = RGBColor(0x58, 0xA6, 0xFF)
GREEN       = RGBColor(0x56, 0xD3, 0x64)
YELLOW      = RGBColor(0xE3, 0xB3, 0x41)
RED         = RGBColor(0xF8, 0x51, 0x49)
WHITE       = RGBColor(0xF0, 0xF6, 0xFC)
GRAY        = RGBColor(0x8B, 0x94, 0x9E)
LIGHT_TEXT  = RGBColor(0xC9, 0xD1, 0xD9)
BORDER      = RGBColor(0x30, 0x36, 0x3D)

BLANK_LAYOUT = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)


def bg(slide, color=DARK_BG):
    bg_shape = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()
    return bg_shape


def box(slide, l, t, w, h, fill=CARD_BG, border=BORDER):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.5)
    return shape


def label(slide, text, l, t, w, h,
          size=14, bold=False, color=LIGHT_TEXT,
          align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def multiline(slide, lines, l, t, w, h, size=13, color=LIGHT_TEXT, spacing=1.15):
    """lines = list of (text, bold, color_or_None)"""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, clr) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = clr if clr else color
    return txb


def code_box(slide, code_text, l, t, w, h, size=11):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x08, 0x0D, 0x13)
    shape.line.color.rgb = RGBColor(0x23, 0x86, 0x36)
    shape.line.width = Pt(1)
    txb = slide.shapes.add_textbox(
        Inches(l + 0.12), Inches(t + 0.1),
        Inches(w - 0.24), Inches(h - 0.2)
    )
    txb.word_wrap = False
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in code_text.strip().split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(0x79, 0xC0, 0xFF)
        run.font.name = "Courier New"
    return txb


def slide_num(slide, n, total=12):
    label(slide, f"{n} / {total}", 12.5, 7.1, 0.7, 0.3,
          size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def tag_label(slide, text, l=0.4, t=0.25):
    label(slide, text, l, t, 5, 0.35, size=10, bold=True, color=BLUE)


def title_text(slide, text, l=0.4, t=0.65, w=12.5, size=32):
    label(slide, text, l, t, w, 1.2, size=size, bold=True, color=WHITE)


def divider(slide, t, l=0.4, w=12.5):
    line = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); slide_num(s, 1)

label(s, "DSCI 551  ·  Spring 2026", 0.4, 0.3, 8, 0.4,
      size=11, bold=True, color=BLUE)

label(s, "Clinical Diagnosis and\nTreatment Records System",
      0.4, 0.8, 10, 1.8, size=36, bold=True, color=WHITE)

label(s, "Exploring PostgreSQL Internals: B-Tree Indexing & MVCC",
      0.4, 2.7, 10, 0.5, size=16, color=GRAY)

divider(s, 3.4)

label(s, "Diana Rulanova  ·  Eliza Loke",
      0.4, 3.6, 8, 0.45, size=15, bold=True, color=BLUE)
label(s, "University of Southern California  ·  DSCI 551 Data Management",
      0.4, 4.1, 10, 0.4, size=12, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Motivation
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); slide_num(s, 2)
tag_label(s, "INTRODUCTION")
title_text(s, "Motivation & Goals")
divider(s, 1.75)

# Left card
box(s, 0.4, 2.0, 5.9, 4.8)
label(s, "Why PostgreSQL?", 0.65, 2.15, 5.5, 0.4,
      size=12, bold=True, color=BLUE)
multiline(s, [
    ("•  Mature, production-grade RDBMS", False, None),
    ("•  Strong support for indexing & transactions", False, None),
    ("•  Internals are observable via EXPLAIN ANALYZE", False, None),
    ("•  MVCC built into the core engine", False, None),
    ("•  Real-world adoption in healthcare systems", False, None),
], 0.65, 2.55, 5.5, 3.5, size=13)

# Right card
box(s, 6.7, 2.0, 6.2, 4.8)
label(s, "Project Goals", 6.95, 2.15, 5.8, 0.4,
      size=12, bold=True, color=BLUE)
multiline(s, [
    ("•  Analyze B-tree index behavior in real queries", False, None),
    ("•  Compare Index Scan vs Sequential Scan", False, None),
    ("•  Demonstrate MVCC with concurrent transactions", False, None),
    ("•  Show row-level locking with SELECT FOR UPDATE", False, None),
    ("•  Map each app operation to a DB internal", False, None),
], 6.95, 2.55, 5.8, 3.5, size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Schema
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); slide_num(s, 3)
tag_label(s, "APPLICATION DESIGN")
title_text(s, "Database Schema")
divider(s, 1.75)

# Entity boxes
entities = [
    ("PATIENTS", "patient_id PK\nfull_name · dob\ncontact", 0.3),
    ("ENCOUNTERS", "encounter_id PK\npatient_id FK\nprovider_id FK\nencounter_date · complaint", 3.1),
    ("DIAGNOSES", "diagnosis_id PK\nencounter_id FK\ndiagnosis_code\ndescription", 6.2),
    ("TREATMENTS", "treatment_id PK\ndiagnosis_id FK\ntreatment_plan\noutcome", 9.3),
]
for name, fields, x in entities:
    box(s, x, 2.0, 2.7, 2.4, fill=RGBColor(0x1F, 0x26, 0x2D))
    label(s, name, x+0.1, 2.1, 2.5, 0.4, size=11, bold=True, color=BLUE)
    label(s, fields, x+0.1, 2.55, 2.5, 1.8, size=10, color=LIGHT_TEXT)

# Arrows between entities
for x in [3.0, 6.1, 9.2]:
    label(s, "→", x - 0.1, 2.9, 0.3, 0.4, size=18, color=GRAY, align=PP_ALIGN.CENTER)

# PROVIDERS
box(s, 0.3, 4.7, 2.7, 2.0, fill=RGBColor(0x1F, 0x26, 0x2D))
label(s, "PROVIDERS", 0.4, 4.8, 2.5, 0.4, size=11, bold=True, color=BLUE)
label(s, "provider_id PK\nfull_name · specialty\ncontact", 0.4, 5.2, 2.5, 1.4, size=10, color=LIGHT_TEXT)
label(s, "↗", 3.05, 4.5, 0.3, 0.5, size=16, color=GRAY)

# Dataset & Indexes info
box(s, 3.1, 4.7, 9.8, 2.1)
label(s, "Dataset & Indexes", 3.3, 4.82, 9.4, 0.4, size=12, bold=True, color=BLUE)
multiline(s, [
    ("Dataset:   1,000 patients  ·  50 providers  ·  10,000 encounters  ·  ~20,000 diagnoses  ·  ~30,000 treatments", False, LIGHT_TEXT),
    ("", False, None),
    ("idx_diagnoses_code              →  B-tree on diagnoses(diagnosis_code)", False, GREEN),
    ("idx_encounters_patient_date  →  B-tree on encounters(patient_id, encounter_date)", False, GREEN),
], 3.3, 5.2, 9.4, 1.5, size=12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – B-Tree Theory
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); slide_num(s, 4)
tag_label(s, "INTERNAL FOCUS  ·  DIANA RULANOVA")
title_text(s, "B-Tree Indexing in PostgreSQL")
divider(s, 1.75)

box(s, 0.4, 2.0, 6.0, 5.2)
label(s, "How B-Tree Works", 0.65, 2.1, 5.6, 0.4, size=12, bold=True, color=BLUE)
multiline(s, [
    ("•  Default index type in PostgreSQL", False, None),
    ("•  Balanced tree: root → internal → leaf pages", False, None),
    ("•  Leaf pages store (key, heap TID) pairs", False, None),
    ("•  Lookup cost:  O(log n)  vs  O(n) for sequential scan", False, None),
    ("", False, None),
    ("Planner strategies:", True, WHITE),
    ("   Index Scan          few matching rows, pointer-by-pointer", False, None),
    ("   Bitmap Index Scan  many rows, batch heap read", False, None),
    ("   Sequential Scan    when index won't help (e.g. aggregations)", False, None),
], 0.65, 2.55, 5.6, 4.5, size=13)

box(s, 6.8, 2.0, 6.1, 5.2)
label(s, "Composite Index Benefit", 7.05, 2.1, 5.7, 0.4, size=12, bold=True, color=BLUE)
label(s, "Index on  (patient_id, encounter_date):", 7.05, 2.55, 5.7, 0.4, size=13, color=LIGHT_TEXT)
code_box(s, "WHERE  patient_id = 50\nORDER  BY encounter_date DESC", 7.05, 3.05, 5.7, 0.85)
multiline(s, [
    ("•  Left prefix  →  filters by patient_id", False, None),
    ("•  Right column  →  data already sorted by date", False, None),
    ("•  No extra Sort node in the query plan", True, GREEN),
    ("", False, None),
    ("This 'skip sort' optimization eliminates", False, None),
    ("a separate sorting step — saving both", False, None),
    ("time and memory for large result sets.", False, None),
], 7.05, 4.0, 5.7, 3.0, size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Operation 1
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); slide_num(s, 5)
tag_label(s, "DEMO  ·  OPERATION 1")
title_text(s, "Diagnosis Lookup by Code")
divider(s, 1.75)

# Left
box(s, 0.4, 2.0, 6.0, 5.2)
label(s, "What the App Does", 0.65, 2.1, 5.6, 0.4, size=12, bold=True, color=BLUE)
label(s, "Clinician searches diagnoses by ICD-style code (e.g. G43 = Migraine).", 0.65, 2.55, 5.6, 0.55, size=13, color=LIGHT_TEXT)
code_box(s, "SELECT diagnosis_id,\n       diagnosis_code,\n       description\nFROM   diagnoses\nWHERE  diagnosis_code = 'G43'", 0.65, 3.15, 5.6, 1.5)
label(s, "Why it matters:", 0.65, 4.75, 5.6, 0.35, size=12, bold=True, color=WHITE)
label(s, "Index reads only 163 pages.\nSequential scan reads all 19,885 rows.\n→  2.6× faster with the index.", 0.65, 5.1, 5.6, 1.0, size=13, color=LIGHT_TEXT)

# Right
box(s, 6.8, 2.0, 6.1, 2.45, fill=RGBColor(0x0A, 0x14, 0x0A))
label(s, "WITH index  (Bitmap Index Scan)", 7.05, 2.1, 5.7, 0.35, size=11, bold=True, color=GREEN)
code_box(s, "Bitmap Heap Scan on diagnoses\n  Recheck Cond: diagnosis_code='G43'\n  Heap Blocks: exact=161\n  -> Bitmap Index Scan on idx_diagnoses_code\nExecution Time: 0.69 ms", 7.05, 2.45, 5.7, 1.85, size=10)

box(s, 6.8, 4.6, 6.1, 2.55, fill=RGBColor(0x14, 0x08, 0x08))
label(s, "WITHOUT index  (Sequential Scan)", 7.05, 4.7, 5.7, 0.35, size=11, bold=True, color=RED)
code_box(s, "Seq Scan on diagnoses\n  Filter: diagnosis_code = 'G43'\n  Rows Removed by Filter: 18,901\nExecution Time: 1.80 ms", 7.05, 5.05, 5.7, 1.85, size=10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Operation 2
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); slide_num(s, 6)
tag_label(s, "DEMO  ·  OPERATION 2")
title_text(s, "Patient Encounter History")
divider(s, 1.75)

box(s, 0.4, 2.0, 6.0, 5.2)
label(s, "What the App Does", 0.65, 2.1, 5.6, 0.4, size=12, bold=True, color=BLUE)
label(s, "Retrieve a patient's visit history, most recent first.", 0.65, 2.55, 5.6, 0.45, size=13, color=LIGHT_TEXT)
code_box(s, "SELECT encounter_id, encounter_date,\n       chief_complaint\nFROM   encounters\nWHERE  patient_id = 50\nORDER  BY encounter_date DESC\nLIMIT  10", 0.65, 3.05, 5.6, 1.7)
label(s, "Key insight:", 0.65, 4.85, 5.6, 0.35, size=12, bold=True, color=WHITE)
label(s, "Composite indexes support both filtering\n(WHERE) and ordering (ORDER BY)\nin a single index traversal.", 0.65, 5.2, 5.6, 1.0, size=13, color=LIGHT_TEXT)

box(s, 6.8, 2.0, 6.1, 3.0, fill=RGBColor(0x0A, 0x14, 0x0A))
label(s, "EXPLAIN ANALYZE output", 7.05, 2.1, 5.7, 0.35, size=11, bold=True, color=GREEN)
code_box(s, "Index Scan Backward\n  using idx_encounters_patient_date\n  Index Cond: (patient_id = 50)\n  Buffers: shared hit=12\nExecution Time: 0.077 ms\n\n→  NO Sort node in the plan!", 7.05, 2.45, 5.7, 2.35, size=10)

box(s, 6.8, 5.1, 6.1, 2.1)
label(s, "Why No Sort Node?", 7.05, 5.2, 5.7, 0.35, size=12, bold=True, color=BLUE)
multiline(s, [
    ("•  Index stores rows in (patient_id, encounter_date) order", False, None),
    ("•  Backward scan = descending date for free", False, None),
    ("•  PostgreSQL reads leaf pages — already sorted", False, None),
], 7.05, 5.6, 5.7, 1.4, size=12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – MVCC Theory
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); slide_num(s, 7)
tag_label(s, "INTERNAL FOCUS  ·  ELIZA LOKE")
title_text(s, "MVCC in PostgreSQL")
divider(s, 1.75)

box(s, 0.4, 2.0, 6.0, 5.2)
label(s, "How MVCC Works", 0.65, 2.1, 5.6, 0.4, size=12, bold=True, color=BLUE)
multiline(s, [
    ("•  Every row version stores xmin (created by) and xmax (deleted by)", False, None),
    ("•  UPDATE creates a new heap tuple — old one stays", False, None),
    ("•  Each transaction sees a consistent snapshot", False, None),
    ("•  Readers never block writers", True, GREEN),
    ("•  Writers never block readers", True, GREEN),
    ("", False, None),
    ("READ COMMITTED isolation (default):", True, WHITE),
    ("•  Each SQL statement sees the latest committed data", False, None),
    ("•  No dirty reads — uncommitted changes are invisible", False, None),
    ("•  New snapshot taken at each statement start", False, None),
], 0.65, 2.55, 5.6, 4.5, size=13)

box(s, 6.8, 2.0, 6.1, 5.2)
label(s, "Row-Level Locking", 7.05, 2.1, 5.7, 0.4, size=12, bold=True, color=BLUE)
multiline(s, [
    ("SELECT FOR UPDATE acquires an exclusive row lock:", False, LIGHT_TEXT),
    ("", False, None),
    ("•  Concurrent writers must wait for the lock", False, None),
    ("•  Readers are NOT blocked (still use MVCC)", False, None),
    ("•  Lock is released on COMMIT or ROLLBACK", False, None),
    ("", False, None),
    ("MVCC vs Locking:", True, WHITE),
    ("", False, None),
    ("MVCC          →  read/write concurrency", False, GREEN),
    ("SELECT FOR UPDATE  →  write/write serialization", False, YELLOW),
    ("", False, None),
    ("Together they handle all concurrency scenarios.", False, LIGHT_TEXT),
], 7.05, 2.55, 5.7, 4.5, size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Operation 3 MVCC
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); slide_num(s, 8)
tag_label(s, "DEMO  ·  OPERATION 3")
title_text(s, "Concurrent Update — MVCC")
divider(s, 1.75)

box(s, 0.4, 2.0, 5.5, 5.2)
label(s, "Timeline", 0.65, 2.1, 5.1, 0.4, size=12, bold=True, color=BLUE)

rows_data = [
    ("Time", "Session 1", "Session 2", WHITE, True),
    ("t1", "UPDATE outcome='Improved'\n(not committed)", "—", LIGHT_TEXT, False),
    ("t2", "… waiting …", "READ → sees 'Improving'\n(old committed value)", LIGHT_TEXT, False),
    ("t3", "COMMIT ✓", "—", LIGHT_TEXT, False),
    ("t4", "—", "READ → sees 'Improved'\n(new value)", LIGHT_TEXT, False),
]
y = 2.6
for row in rows_data:
    t_val, s1_val, s2_val, clr, hdr_row = row
    is_hdr = hdr_row
    label(s, t_val, 0.55, y, 0.7, 0.55, size=11, bold=is_hdr, color=BLUE if is_hdr else GRAY)
    label(s, s1_val, 1.3, y, 2.1, 0.55, size=11, bold=is_hdr, color=WHITE if is_hdr else LIGHT_TEXT)
    label(s, s2_val, 3.5, y, 2.2, 0.55, size=11, bold=is_hdr, color=WHITE if is_hdr else GREEN)
    y += 0.62

box(s, 6.2, 2.0, 6.7, 5.2)
label(s, "What PostgreSQL Does Internally", 6.45, 2.1, 6.3, 0.4, size=12, bold=True, color=BLUE)
multiline(s, [
    ("Session 1's UPDATE creates a new heap tuple:", False, LIGHT_TEXT),
    ("   xmin = Session 1's transaction ID", False, LIGHT_TEXT),
    ("   Old tuple's xmax = Session 1's xid", False, LIGHT_TEXT),
    ("", False, None),
    ("Session 2 checks tuple visibility:", False, LIGHT_TEXT),
    ("   Session 1's xid is NOT in committed set", False, LIGHT_TEXT),
    ("   → Returns the old committed version", True, GREEN),
    ("", False, None),
    ("After Session 1 COMMIT:", False, LIGHT_TEXT),
    ("   Session 1's xid is now committed", False, LIGHT_TEXT),
    ("   → New version becomes visible", True, YELLOW),
    ("", False, None),
    ("Why it matters:", True, WHITE),
    ("No locking needed between reader and writer.", False, LIGHT_TEXT),
    ("Critical in healthcare: reads must never", False, LIGHT_TEXT),
    ("stall due to an in-progress update.", False, LIGHT_TEXT),
], 6.45, 2.55, 6.3, 4.5, size=12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Operation 4 Locking
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); slide_num(s, 9)
tag_label(s, "DEMO  ·  OPERATION 4")
title_text(s, "Row-Level Locking")
divider(s, 1.75)

box(s, 0.4, 2.0, 6.0, 5.2)
label(s, "Scenario", 0.65, 2.1, 5.6, 0.4, size=12, bold=True, color=BLUE)
label(s, "A clinician locks a record before writing to prevent conflicting updates.", 0.65, 2.5, 5.6, 0.6, size=13, color=LIGHT_TEXT)
code_box(s, "-- Session 1\nSELECT outcome FROM treatments\nWHERE  treatment_id = 2\nFOR UPDATE;\n\n-- (holds lock for 4 seconds)\n\nUPDATE treatments\n  SET  outcome = 'Recovered'\n  WHERE treatment_id = 2;\nCOMMIT;", 0.65, 3.15, 5.6, 2.65, size=10)

box(s, 6.8, 2.0, 6.1, 5.2)
label(s, "What Happens", 7.05, 2.1, 5.7, 0.4, size=12, bold=True, color=BLUE)
multiline(s, [
    ("1.  Session 1 acquires exclusive row lock", False, None),
    ("2.  Session 2 attempts UPDATE on same row", False, None),
    ("3.  Session 2 BLOCKS — waits 3.5 seconds", True, RED),
    ("4.  Session 1 commits → lock released", False, None),
    ("5.  Session 2 unblocks and succeeds", True, GREEN),
    ("", False, None),
    ("MVCC vs SELECT FOR UPDATE:", True, WHITE),
    ("", False, None),
    ("MVCC handles read/write concurrency.", False, GREEN),
    ("Readers see old version, no blocking.", False, LIGHT_TEXT),
    ("", False, None),
    ("SELECT FOR UPDATE handles write/write.", False, YELLOW),
    ("Prevents last-write-wins conflicts.", False, LIGHT_TEXT),
    ("", False, None),
    ("In healthcare: two nurses cannot", False, LIGHT_TEXT),
    ("overwrite each other's treatment notes.", False, LIGHT_TEXT),
], 7.05, 2.55, 5.7, 4.5, size=12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Operation 5 Aggregation
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); slide_num(s, 10)
tag_label(s, "DEMO  ·  OPERATION 5")
title_text(s, "Monthly Diagnosis Summary")
divider(s, 1.75)

box(s, 0.4, 2.0, 6.0, 5.2)
label(s, "What the App Does", 0.65, 2.1, 5.6, 0.4, size=12, bold=True, color=BLUE)
label(s, "Administrator generates monthly report of diagnosis counts.", 0.65, 2.5, 5.6, 0.5, size=13, color=LIGHT_TEXT)
code_box(s, "SELECT\n  DATE_TRUNC('month', e.encounter_date)\n    AS month,\n  d.diagnosis_code,\n  COUNT(*) AS diagnosis_count\nFROM diagnoses d\nJOIN encounters e\n     ON d.encounter_id = e.encounter_id\nGROUP BY 1, 2\nORDER BY month DESC,\n         diagnosis_count DESC", 0.65, 3.05, 5.6, 3.0, size=10)

box(s, 6.8, 2.0, 6.1, 5.2)
label(s, "Query Plan", 7.05, 2.1, 5.7, 0.4, size=12, bold=True, color=BLUE)
code_box(s, "Hash Join\n  -> Seq Scan on diagnoses\n       (19,885 rows scanned)\n  -> Hash / Seq Scan on encounters\n       (10,000 rows scanned)\nHashAggregate\n  Group Key: month, diagnosis_code\nExecution Time: 25.9 ms", 7.05, 2.5, 5.7, 2.2, size=10)

label(s, "Why Sequential Scan here?", 7.05, 4.8, 5.7, 0.4, size=12, bold=True, color=BLUE)
multiline(s, [
    ("•  Query must touch ALL rows to compute COUNT(*)", False, None),
    ("•  Index scan adds overhead with no benefit", False, None),
    ("•  Planner correctly chooses Seq Scan + Hash Join", True, GREEN),
    ("•  Contrasts with Ops 1 & 2 where queries are selective", False, None),
    ("", False, None),
    ("The planner is smart — it uses indexes only", False, LIGHT_TEXT),
    ("when they actually help.", False, LIGHT_TEXT),
], 7.05, 5.2, 5.7, 1.95, size=12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Comparison
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); slide_num(s, 11)
tag_label(s, "ANALYSIS")
title_text(s, "PostgreSQL vs MySQL vs MongoDB")
divider(s, 1.75)

headers = ["Feature", "PostgreSQL", "MySQL (InnoDB)", "MongoDB"]
col_x   = [0.35, 3.3, 6.65, 10.0]
col_w   = [2.85, 3.25, 3.25, 3.1]

# header row
for i, h in enumerate(headers):
    box(s, col_x[i], 2.0, col_w[i], 0.5, fill=RGBColor(0x1F, 0x6F, 0xEB), border=RGBColor(0x1F, 0x6F, 0xEB))
    label(s, h, col_x[i]+0.1, 2.05, col_w[i]-0.1, 0.4, size=12, bold=True, color=WHITE)

table_data = [
    ("Primary index",    "Heap + secondary B-tree",   "Clustered B-tree (PK)",     "B-tree on _id"),
    ("Concurrency",      "MVCC — no read locks",       "MVCC + gap locks",           "Document-level locking"),
    ("Composite index",  "Skip sort, prefix scan",     "Similar support",            "Compound indexes"),
    ("EXPLAIN",          "EXPLAIN ANALYZE + BUFFERS",  "EXPLAIN FORMAT=JSON",        "explain('executionStats')"),
    ("Schema",           "Strict relational",          "Strict relational",          "Flexible document"),
    ("Write model",      "Heap append + WAL",          "In-place update + redo log", "WiredTiger journal"),
]

highlight_col = [1, 1, 1, 1, 0, 0]  # which column to highlight per row

for r_idx, row in enumerate(table_data):
    y = 2.55 + r_idx * 0.62
    fill = RGBColor(0x0D, 0x11, 0x17) if r_idx % 2 == 0 else RGBColor(0x12, 0x17, 0x1E)
    for c_idx, cell in enumerate(row):
        box(s, col_x[c_idx], y, col_w[c_idx], 0.6, fill=fill)
        clr = GREEN if c_idx == 1 and c_idx != 0 else LIGHT_TEXT
        if c_idx == 0:
            clr = GRAY
        label(s, cell, col_x[c_idx]+0.1, y+0.07, col_w[c_idx]-0.1, 0.5,
              size=11, color=clr)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – Conclusion
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); slide_num(s, 12)
tag_label(s, "CONCLUSION")
title_text(s, "Summary & Lessons Learned")
divider(s, 1.75)

cards = [
    ("B-Tree Indexing", [
        "2.6× faster lookup with index",
        "Bitmap scan for high-selectivity queries",
        "Composite index eliminates Sort node",
        "Seq Scan wins for full-table aggregations",
        "Planner decides automatically based on cost",
    ], BLUE),
    ("MVCC", [
        "No dirty reads under READ COMMITTED",
        "Readers and writers don't block each other",
        "Row versioning via xmin/xmax in the heap",
        "SELECT FOR UPDATE for write serialization",
        "Both mechanisms work together",
    ], GREEN),
    ("Key Takeaway", [
        "Every app operation maps to a DB internal",
        "Index design directly affects the query plan",
        "Dataset size determines if index is used",
        "MVCC is why PostgreSQL scales under concurrency",
        "EXPLAIN ANALYZE makes internals observable",
    ], YELLOW),
]

for i, (title, points, color) in enumerate(cards):
    x = 0.4 + i * 4.3
    box(s, x, 2.05, 4.1, 5.1)
    label(s, title, x+0.2, 2.15, 3.8, 0.45, size=13, bold=True, color=color)
    lines = [("• " + p, False, None) for p in points]
    multiline(s, lines, x+0.2, 2.65, 3.8, 4.3, size=12)

divider(s, 7.1)
label(s, "Diana Rulanova  ·  Eliza Loke  ·  DSCI 551 Spring 2026",
      0.4, 7.15, 12.5, 0.35, size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "DSCI551_Clinical_Records_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
